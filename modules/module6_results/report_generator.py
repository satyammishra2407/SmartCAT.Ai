"""PPTX / PDF summary reports."""
from __future__ import annotations

from pathlib import Path
from typing import Any

import pandas as pd
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from modules.module6_results.chart_generator import plot_oep_aep
from modules.module6_results.metrics_calculator import metrics_from_ep
from modules.module6_results.output_parser import load_ep_curve
from smartcat_logging import get_logger

logger = get_logger("module6.report")


class ResultsInterpretationEngine:
    def __init__(self, output_dir: Path | None = None):
        self.output_dir = Path(output_dir) if output_dir else Path("output/reports")

    def run(self, model_output_csv: Path, stem: str = "results_summary") -> dict[str, Path]:
        self.output_dir.mkdir(parents=True, exist_ok=True)
        df = load_ep_curve(model_output_csv)
        m = metrics_from_ep(df)
        rp = m["return_period_table"]
        stats = m["stats"]

        png = self.output_dir / f"{stem}_oep_aep.png"
        plot_oep_aep(rp, png)

        pdf_path = self.output_dir / f"{stem}_report.pdf"
        self._write_pdf(pdf_path, stats, rp)

        pptx_path = self.output_dir / f"{stem}_report.pptx"
        self._write_pptx(pptx_path, stats, rp, png)

        rp.to_csv(self.output_dir / f"{stem}_return_periods.csv", index=False)

        return {"pdf": pdf_path, "pptx": pptx_path, "chart": png, "rp_csv": self.output_dir / f"{stem}_return_periods.csv"}

    def _write_pdf(self, path: Path, stats: dict[str, Any], rp: pd.DataFrame) -> None:
        c = canvas.Canvas(str(path), pagesize=letter)
        width, height = letter
        y = height - 72
        c.setFont("Helvetica-Bold", 14)
        c.drawString(72, y, "SmartCAT.AI — Results Summary")
        y -= 28
        c.setFont("Helvetica", 11)
        summary = (
            f"Average Annual Loss (sample metric): {stats.get('aal')}\n"
            f"Std Dev: {stats.get('std')} | CV: {stats.get('cv')}\n\n"
            "Return period losses are interpolated from supplied EP data. "
            "Validate against vendor documentation before client distribution."
        )
        for line in summary.split("\n"):
            c.drawString(72, y, line[:120])
            y -= 14
        y -= 10
        c.drawString(72, y, "Return period table (preview):")
        y -= 16
        for _, row in rp.head(10).iterrows():
            c.drawString(72, y, str(dict(row)))
            y -= 14
            if y < 100:
                c.showPage()
                y = height - 72
        c.save()

    def _write_pptx(self, path: Path, stats: dict[str, Any], rp: pd.DataFrame, chart_png: Path) -> None:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            logger.warning("python-pptx not available")
            return

        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "SmartCAT.AI Results"

        left = Inches(0.5)
        top = Inches(1.2)
        box = slide.shapes.add_textbox(left, top, Inches(9), Inches(1.5))
        tf = box.text_frame
        tf.text = (
            f"AAL (heuristic): {stats.get('aal')}\n"
            f"CV: {stats.get('cv')} — see CSV/PDF for full return period table."
        )

        if chart_png.exists():
            slide.shapes.add_picture(str(chart_png), Inches(1), Inches(2.8), width=Inches(8))

        prs.save(str(path))
